from pathlib import Path
import re
from pptx import Presentation
from PIL import Image
import os
from slugify import slugify

class PPTXConverter:
    def __init__(self, input_file):
        self.pptx = Presentation(input_file)
        self.input_path = Path(input_file).resolve()
        
        # Extract chapter number from input filename (e.g., "479-wk01.pptx" or "479-ch01-1_8.pptx" -> "01")
        chapter_match = re.search(r'(?:wk|ch)(\d+)', self.input_path.stem)
        self.chapter = chapter_match.group(1) if chapter_match else "00"
        
        # Setup directory structure relative to slides directory
        self.slides_dir = Path('/Users/ovandef/Documents/GitHub/kin479/slides')
        self.chapter_dir = self.slides_dir / f"ch{self.chapter}" / "source-files"
        self.common_dir = self.slides_dir / "common-all-slides" / "js"
        self.image_dir = self.chapter_dir / "images"
        
        self.setup_directories()

    def setup_directories(self):
        """Create necessary directories."""
        self.chapter_dir.mkdir(exist_ok=True)
        self.common_dir.mkdir(exist_ok=True)
        self.image_dir.mkdir(exist_ok=True)

        # Create common styles.css if it doesn't exist
        styles_file = self.common_dir / "styles.css"
        if not styles_file.exists():
            self.create_styles_file(styles_file)

    def create_styles_file(self, styles_file):
        """Create the common styles CSS file."""
        with open(styles_file, 'w', encoding='utf-8') as f:
            f.write("""
/* Base presentation styles */
.reveal {
    font-size: 40px;
}

/* Title slide customization */
.reveal .title-slide {
    text-align: left;
}

.reveal .title-slide h1 {
    font-size: 2.5em;
    margin-bottom: 0.5em;
}

.reveal .title-slide h2 {
    font-size: 1.8em;
    color: #666;
}

.reveal .title-slide .author {
    margin-top: 2em;
}

.reveal .title-slide .affiliation {
    font-style: italic;
    color: #666;
}

/* Content styles */
.reveal h2 {
    font-size: 1.8em;
    margin-bottom: 0.5em;
}

.reveal ul {
    font-size: 1em;
    line-height: 1.4;
}

.reveal li {
    margin-bottom: 0.5em;
}

/* Image and caption styles */
.reveal img {
    margin: 0 auto;
    display: block;
}

.reveal figcaption {
    font-size: 0.8em;
    color: #666;
    text-align: center;
    margin-top: 0.5em;
}

/* Column layout */
.reveal .slides section {
    height: 100%;
}

.reveal .column {
    font-size: 1em;
}

.reveal .columns {
    display: flex;
    align-items: center;
}
""")

    def get_quarto_yaml(self):
        """Generate YAML front matter for Quarto presentation."""
        yaml = [
            "---",
            f"title: \"Chapter {self.chapter}\"",
            "subtitle: \"Course Introduction\"",
            "",
            "author:",
            "  - name: Ovande Furtado Jr",
            "    title: Assistant Professor",
            "    department: Kinesiology/KIN479-Motor Control",
            "    orcid: 0000-0003-3847-6314",
            "    email: ovandef@csun.edu",
            "    affiliations: California State University, Northridge",
            "",
            "# Document Settings",
            "date: last-modified",
            "logo: \"/images/logo.png\"",
            "",
            "format: ",
            "  revealjs:",
            "    # Theme and Styling",
            "    theme: default",
            "    css: ../common-all-slides/js/styles.css",
            "    width: 1600",
            "    height: 900",
            "    fontsize: 40px",
            "    margin: 0.1",
            "",
            "    # Navigation Controls",
            "    controls: true",
            "    controls-layout: edges",
            "    controls-back-arrows: visible",
            "    controls-position: bottom",
            "    controls-tutorial: true",
            "    navigation-mode: linear",
            "    progress: true",
            "    slide-number: c/t",
            "    show-slide-number: all",
            "",
            "    # Menu Settings",
            "    menu: true",
            "    menu-position: left",
            "    editor: source",
            "---\n"
        ]
        return "\n".join(yaml)

    def extract_text_from_text_frame(self, text_frame, processed_text=None):
        """Extract text from a text frame with bullet points and exact text preservation."""
        if processed_text is None:
            processed_text = set()
            
        text_parts = []
        
        for paragraph in text_frame.paragraphs:
            if not paragraph.text.strip():
                continue
            
            # Preserve exact indentation level and add bullet points
            prefix = "    " * paragraph.level
            if not prefix:  # Top level
                prefix = "- "
            elif paragraph.level > 0:  # Nested levels
                prefix = prefix[:-4] + "    - "  # Replace last indentation with bullet
            
            # Combine all runs in the paragraph preserving exact text
            para_text = []
            for run in paragraph.runs:
                if run.text:
                    # Preserve exact text, only clean up line endings
                    text = run.text
                    text = text.replace('\v', '\n')  # Convert vertical tabs to newlines
                    text = text.replace('\r', '\n')  # Convert carriage returns to newlines
                    para_text.append(text)
            
            if para_text:
                # Keep original text exactly as is
                text = ' '.join(para_text)
                
                # Normalize text for comparison (remove extra spaces, convert to lowercase)
                normalized_text = ' '.join(text.lower().split())
                
                # Only add text if we haven't seen a normalized version before
                if normalized_text not in processed_text:
                    # Only add bullet if text doesn't already start with one
                    if not text.strip().startswith(('-', '•', '*')):
                        final_text = f"{prefix}{text}"
                    else:
                        final_text = f"{prefix[:-2]}{text}"  # Remove the "- " if text has bullet
                    
                    text_parts.append(final_text)
                    processed_text.add(normalized_text)
        
        return text_parts

    def extract_text(self, shape, processed_text=None):
        """Extract text from a shape with bullet points and exact text preservation."""
        if processed_text is None:
            processed_text = set()
            
        text_parts = []
        
        try:
            # Try to get text from shape properties
            if hasattr(shape, "text_frame") and shape.text_frame:
                if shape.text_frame.text.strip():
                    text_parts.extend(self.extract_text_from_text_frame(shape.text_frame, processed_text))
            
            # Try to get text from placeholders
            if hasattr(shape, "placeholder_format"):
                if hasattr(shape.placeholder_format, "text") and shape.placeholder_format.text.strip():
                    text = shape.placeholder_format.text.strip()
                    normalized_text = ' '.join(text.lower().split())
                    if normalized_text not in processed_text:
                        if not text.startswith(('-', '•', '*')):
                            text_parts.append(f"- {text}")
                        else:
                            text_parts.append(text)
                        processed_text.add(normalized_text)
            
            # Try to get text from table cells
            if hasattr(shape, "table"):
                for row in shape.table.rows:
                    row_text = []
                    for cell in row.cells:
                        if hasattr(cell, "text_frame") and cell.text_frame.text.strip():
                            text = cell.text_frame.text.strip()
                            normalized_text = ' '.join(text.lower().split())
                            if normalized_text not in processed_text:
                                row_text.append(text)
                                processed_text.add(normalized_text)
                    if row_text:
                        text_parts.append(f"- {' | '.join(row_text)}")
            
            # Try to get text from text box
            if hasattr(shape, "text_box"):
                if shape.text_box.text.strip():
                    for line in shape.text_box.text.strip().split('\n'):
                        if line.strip():
                            normalized_text = ' '.join(line.lower().split())
                            if normalized_text not in processed_text:
                                if not line.strip().startswith(('-', '•', '*')):
                                    text_parts.append(f"- {line.strip()}")
                                else:
                                    text_parts.append(line.strip())
                                processed_text.add(normalized_text)
            
            # Try to get text from shape text (fallback)
            if hasattr(shape, "text") and shape.text.strip():
                # Clean up only line endings
                text = shape.text
                text = text.replace('\v', '\n')  # Convert vertical tabs to newlines
                text = text.replace('\r', '\n')  # Convert carriage returns to newlines
                
                for line in text.split('\n'):
                    if line.strip():
                        normalized_text = ' '.join(line.lower().split())
                        if normalized_text not in processed_text:
                            if not line.strip().startswith(('-', '•', '*')):
                                text_parts.append(f"- {line}")
                            else:
                                text_parts.append(line)
                            processed_text.add(normalized_text)
            
            # Try to get text from notes - preserve exactly as is
            if hasattr(shape, "notes_slide") and shape.notes_slide:
                for note_shape in shape.notes_slide.shapes:
                    if hasattr(note_shape, "text") and note_shape.text:
                        text = note_shape.text
                        if text not in processed_text:
                            text_parts.append(text)
                            processed_text.add(text)
        
        except Exception as e:
            print(f"Warning: Error extracting text from shape: {e}")
        
        return '\n'.join(text_parts)

    def save_shape_as_image(self, shape, slide_number, column_width=None):
        """Save shape as image if it's a picture."""
        try:
            # Handle direct image shapes
            if hasattr(shape, "image"):
                image_bytes = shape.image.blob
                base_name = slugify(shape.name) if shape.name else "image"
                image_filename = f"slide_{slide_number}_{base_name}.png"
                image_path = self.image_dir / image_filename
                
                with open(image_path, 'wb') as f:
                    f.write(image_bytes)
                
                # Add scaling attributes for wider columns
                if column_width and column_width > 50:
                    return f"![](images/{image_filename}){{width=100%}}"
                return f"![](images/{image_filename})"
            
            # Handle embedded pictures
            elif hasattr(shape, "element"):
                # Try to find picture element
                pic = shape.element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/picture}pic')
                if pic is not None:
                    # Find blip element that contains image data
                    blip = pic.find('.//{http://schemas.openxmlformats.org/drawingml/2006/picture}blipFill//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
                    if blip is not None:
                        # Get embedded image data
                        image_rid = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                        if image_rid and hasattr(shape.part, 'related_parts'):
                            image_part = shape.part.related_parts[image_rid]
                            image_bytes = image_part.blob
                            
                            # Generate filename and save
                            base_name = slugify(shape.name) if shape.name else "image"
                            image_filename = f"slide_{slide_number}_{base_name}.png"
                            image_path = self.image_dir / image_filename
                            
                            with open(image_path, 'wb') as f:
                                f.write(image_bytes)
                            
                            # Add scaling attributes for wider columns
                            if column_width and column_width > 50:
                                return f"![](images/{image_filename}){{width=100%}}"
                            return f"![](images/{image_filename})"
                
                # Try to find graphic frame element (for charts, diagrams, etc.)
                graphic = shape.element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}graphic')
                if graphic is not None:
                    # Save the shape as an image using PIL
                    try:
                        from PIL import Image
                        import io
                        
                        # Get shape dimensions
                        width = shape.width
                        height = shape.height
                        
                        # Create a new image with the shape dimensions
                        image = Image.new('RGB', (width, height), 'white')
                        
                        # Save the image
                        base_name = slugify(shape.name) if shape.name else "diagram"
                        image_filename = f"slide_{slide_number}_{base_name}.png"
                        image_path = self.image_dir / image_filename
                        
                        image.save(image_path, 'PNG')
                        
                        return f"![](images/{image_filename})"
                    except Exception as e:
                        print(f"Warning: Could not save graphic as image: {e}")
        
        except Exception as e:
            print(f"Warning: Could not save image from shape: {e}")
        
        return ""

    def convert_slide(self, slide, slide_number):
        """Convert a single slide to markdown."""
        if not self.has_content(slide):
            return ""
            
        md_content = []
        
        # Add slide title exactly as is
        title = self.get_slide_title(slide, slide_number)
        md_content.append(f"## {title}\n")
        
        # Collect text and images separately
        text_content = []
        image_content = []
        note_content = []
        
        # Keep track of processed text to avoid duplicates
        processed_text = set()
        
        # Process shapes in reading order (top to bottom, left to right)
        try:
            shapes = sorted(slide.shapes, key=lambda s: (s.top, s.left) if hasattr(s, 'top') and hasattr(s, 'left') else (0, 0))
        except Exception:
            shapes = slide.shapes
        
        # First pass: collect all text and images
        for shape in shapes:
            # Skip if it's the title we already processed
            if shape == slide.shapes.title:
                continue
            
            # Extract text exactly as is
            text = self.extract_text(shape, processed_text)
            if text:
                text_content.extend(text.split('\n'))
            
            # Handle images
            image_md = self.save_shape_as_image(shape, slide_number, column_width=65)  # Pass column width
            if image_md:
                image_content.append(image_md)
        
        # Always use two-column layout if there are images
        if image_content:
            md_content.append("::::: columns\n")
            
            # Text column (65%)
            md_content.append("::: {.column width=\"65%\"}")
            if text_content:
                md_content.extend(text_content)
            md_content.append(":::\n")
            
            # Image column (35%)
            md_content.append("::: {.column width=\"35%\"}")
            md_content.extend(image_content)
            md_content.append(":::")
            
            md_content.append(":::::\n")
        else:
            # If we only have text, display it normally
            if text_content:
                md_content.extend(text_content)
                md_content.append("")
        
        # Extract speaker notes from the slide - preserve exactly as is
        if hasattr(slide, 'notes_slide') and slide.notes_slide:
            for shape in slide.notes_slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    notes_text = shape.text_frame.text
                    if notes_text:
                        note_content.append(notes_text.strip())  # Strip any trailing whitespace
        
        # Add speaker notes at the bottom if any
        if note_content:
            md_content.append("\n::: {.notes}")
            md_content.extend(note_content)
            md_content.append(":::\n")
        
        return "\n".join(md_content)

    def get_slide_title(self, slide, slide_number):
        """Get slide title or generate a default one."""
        if slide.shapes.title and slide.shapes.title.text.strip():
            return slide.shapes.title.text.strip()
        
        # Try to find first text content for default title
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Use first line of text, limited to 50 characters
                default_title = shape.text.strip().split('\n')[0][:50]
                return f"{default_title}..."
        
        # If no text found, use generic title
        return f"Slide {slide_number}"

    def convert(self):
        """Convert the entire presentation to Quarto markdown."""
        md_content = []
        
        # Add Quarto YAML front matter
        md_content.append(self.get_quarto_yaml())
        
        # Convert each slide
        slides = [slide for slide in self.pptx.slides if self.has_content(slide)]
        for i, slide in enumerate(slides, 1):
            if i > 1:
                md_content.append("\n---\n")  # Use standard Quarto slide separator
            md_content.append(self.convert_slide(slide, i))

        # Save the markdown file in the chapter directory
        output_file = self.chapter_dir / f"{self.input_path.stem}.qmd"
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write("\n".join(md_content))
        
        return output_file

    def has_content(self, slide):
        """Check if a slide has any content (text or images)."""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                return True
            if hasattr(shape, "image"):
                return True
        return False


def convert_pptx(input_file):
    """Convert a PowerPoint file to Quarto markdown."""
    converter = PPTXConverter(input_file)
    output_file = converter.convert()
    print(f"Conversion complete. Output saved to: {output_file}")


if __name__ == "__main__":
    import sys
    if len(sys.argv) != 2:
        print("Usage: python pptx_converter.py <input_file>")
        sys.exit(1)
    
    # Get absolute path of input file
    input_file = Path(sys.argv[1]).resolve()
    if not input_file.exists():
        print(f"Error: File {input_file} does not exist")
        sys.exit(1)
    
    convert_pptx(input_file)
